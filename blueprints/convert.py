"""
Conversion Blueprint
  POST /api/to-jpg    — body: file, dpi (72|150|300)
  POST /api/to-png    — body: file, dpi
  POST /api/to-webp   — body: file, quality (1-100)
  POST /api/to-excel  — body: file
  POST /api/to-word   — body: file
  POST /api/to-ppt    — body: file
"""
import os
import io
import base64
import zipfile
import json
import tempfile
from flask import Blueprint, request, jsonify
import fitz
from .utils import save_upload, json_ok, json_err, unique_path, download_url

convert_bp = Blueprint('convert', __name__)

# ─────────────────────────── helpers ────────────────────────────────────────

def _pdf_to_images(src: str, fmt: str, dpi: int = 150, quality: int = 85) -> tuple[list[bytes], int]:
    """Render every page of a PDF to image bytes. Returns (list_of_bytes, page_count)."""
    doc = fitz.open(src)
    images = []
    zoom = dpi / 72.0
    mat = fitz.Matrix(zoom, zoom)

    for page in doc:
        pix = page.get_pixmap(matrix=mat, alpha=(fmt == 'png'))
        if fmt == 'jpg':
            data = pix.tobytes('jpeg', jpg_quality=quality)
        elif fmt == 'png':
            data = pix.tobytes('png')
        elif fmt == 'webp':
            # PyMuPDF does not natively produce WebP; render PNG then convert via Pillow
            from PIL import Image
            png_data = pix.tobytes('png')
            im = Image.open(io.BytesIO(png_data))
            buf = io.BytesIO()
            im.save(buf, 'WEBP', quality=quality)
            data = buf.getvalue()
        else:
            data = pix.tobytes('png')
        images.append(data)

    page_count = doc.page_count
    doc.close()
    return images, page_count


def _images_to_zip(images: list[bytes], fmt: str, base_name: str) -> bytes:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zf:
        for i, img in enumerate(images, 1):
            zf.writestr(f'{base_name}_page{i:03d}.{fmt}', img)
    return buf.getvalue()


# ───────────────────────────── image endpoints ───────────────────────────────

def _image_endpoint(fmt: str):
    if 'file' not in request.files:
        return jsonify(*json_err('No file provided'))

    src, meta = save_upload(request.files['file'])
    if not src:
        return jsonify(*json_err(meta.get('error', 'Upload failed')))

    try:
        dpi     = int(request.form.get('dpi', 150))
        quality = int(request.form.get('quality', 85))
        dpi     = max(36, min(600, dpi))
        quality = max(1, min(100, quality))
    except (TypeError, ValueError):
        dpi, quality = 150, 85

    try:
        images, page_count = _pdf_to_images(src, fmt, dpi, quality)
    except Exception as exc:
        return jsonify(*json_err(f'Conversion failed: {exc}'))
    finally:
        try:
            os.remove(src)
        except OSError:
            pass

    if page_count == 1:
        dest = unique_path(fmt)
        with open(dest, 'wb') as f:
            f.write(images[0])
        return jsonify(**json_ok({
            'download_url': download_url(dest),
            'type': fmt,
            'page_count': 1,
            'output_size': os.path.getsize(dest),
            'warning': meta.get('warning'),
        })[0])
    else:
        zip_bytes = _images_to_zip(images, fmt, 'pages')
        dest = unique_path('zip')
        with open(dest, 'wb') as f:
            f.write(zip_bytes)
        return jsonify(**json_ok({
            'download_url': download_url(dest),
            'type': 'zip',
            'page_count': page_count,
            'output_size': os.path.getsize(dest),
            'warning': meta.get('warning'),
        })[0])


@convert_bp.route('/api/to-jpg', methods=['POST'])
def to_jpg():
    return _image_endpoint('jpg')


@convert_bp.route('/api/to-png', methods=['POST'])
def to_png():
    return _image_endpoint('png')


@convert_bp.route('/api/to-webp', methods=['POST'])
def to_webp():
    return _image_endpoint('webp')


# ───────────────────────────── Excel endpoint ────────────────────────────────

@convert_bp.route('/api/to-excel', methods=['POST'])
def to_excel():
    if 'file' not in request.files:
        return jsonify(*json_err('No file provided'))

    src, meta = save_upload(request.files['file'])
    if not src:
        return jsonify(*json_err(meta.get('error', 'Upload failed')))

    dest = unique_path('xlsx')

    try:
        import pdfplumber
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

        wb = openpyxl.Workbook()
        wb.remove(wb.active)  # remove default sheet

        table_count = 0
        preview_tables = []

        with pdfplumber.open(src) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                tables = page.extract_tables()
                if not tables:
                    # Try extracting plain text into a sheet
                    text = page.extract_text()
                    if text and text.strip():
                        ws = wb.create_sheet(title=f'Page{page_num}_text')
                        for line in text.split('\n'):
                            ws.append([line])
                    continue

                for t_idx, table in enumerate(tables, 1):
                    sheet_name = f'P{page_num}T{t_idx}'
                    ws = wb.create_sheet(title=sheet_name[:31])
                    table_count += 1

                    # Header styling
                    header_fill = PatternFill('solid', fgColor='4F46E5')
                    header_font = Font(bold=True, color='FFFFFF')
                    thin = Side(border_style='thin', color='CCCCCC')
                    border = Border(left=thin, right=thin, top=thin, bottom=thin)

                    for row_i, row in enumerate(table, 1):
                        for col_i, cell_val in enumerate(row or [], 1):
                            cell = ws.cell(row=row_i, column=col_i,
                                           value=str(cell_val) if cell_val is not None else '')
                            cell.border = border
                            cell.alignment = Alignment(wrap_text=True)
                            if row_i == 1:
                                cell.fill = header_fill
                                cell.font = header_font

                    # Auto-fit column widths (capped at 50)
                    for col in ws.columns:
                        max_w = 0
                        for c in col:
                            if c.value:
                                max_w = max(max_w, len(str(c.value)))
                        ws.column_dimensions[col[0].column_letter].width = min(max_w + 2, 50)

                    # Collect preview (first table only, first 5 rows)
                    if len(preview_tables) < 5 and table:
                        preview_tables.append({
                            'sheet': sheet_name,
                            'rows': [
                                [str(c) if c else '' for c in row]
                                for row in table[:5]
                            ],
                        })

        if not wb.sheetnames:
            return jsonify(*json_err(
                'No tables or text found in this PDF. '
                'If it is a scanned image-PDF, OCR is required first.'
            ))

        wb.save(dest)

    except ImportError as exc:
        return jsonify(*json_err(f'Missing library: {exc}. Run: pip install pdfplumber openpyxl'))
    except Exception as exc:
        return jsonify(*json_err(f'Excel conversion failed: {exc}'))
    finally:
        try:
            os.remove(src)
        except OSError:
            pass

    return jsonify(**json_ok({
        'download_url': download_url(dest),
        'table_count': table_count,
        'preview': preview_tables,
        'output_size': os.path.getsize(dest),
        'warning': meta.get('warning'),
    })[0])


# ───────────────────────────── Word endpoint ─────────────────────────────────

@convert_bp.route('/api/to-word', methods=['POST'])
def to_word():
    if 'file' not in request.files:
        return jsonify(*json_err('No file provided'))

    src, meta = save_upload(request.files['file'])
    if not src:
        return jsonify(*json_err(meta.get('error', 'Upload failed')))

    dest = unique_path('docx')

    # Detect if PDF is image-based
    image_based = _is_image_pdf(src)

    try:
        from pdf2docx import Converter
        cv = Converter(src)
        cv.convert(dest, start=0, end=None)
        cv.close()
    except ImportError:
        return jsonify(*json_err('pdf2docx not installed. Run: pip install pdf2docx'))
    except Exception as exc:
        return jsonify(*json_err(f'Word conversion failed: {exc}'))
    finally:
        try:
            os.remove(src)
        except OSError:
            pass

    if not os.path.exists(dest) or os.path.getsize(dest) == 0:
        return jsonify(*json_err('Conversion produced an empty file'))

    warnings = []
    if meta.get('warning'):
        warnings.append(meta['warning'])
    if image_based:
        warnings.append(
            'This PDF appears to be image-based (scanned). '
            'Text extraction quality may be limited. '
            'For best results, use an OCR tool first.'
        )

    return jsonify(**json_ok({
        'download_url': download_url(dest),
        'output_size': os.path.getsize(dest),
        'image_based': image_based,
        'warning': ' | '.join(warnings) if warnings else None,
    })[0])


def _is_image_pdf(path: str) -> bool:
    """Heuristic: if the first 3 pages have no extractable text, assume image-based."""
    try:
        doc = fitz.open(path)
        limit = min(3, doc.page_count)
        total_chars = sum(len(doc[i].get_text()) for i in range(limit))
        doc.close()
        return total_chars < 20
    except Exception:
        return False


# ─────────────────────────── PowerPoint endpoint ─────────────────────────────

@convert_bp.route('/api/to-ppt', methods=['POST'])
def to_ppt():
    if 'file' not in request.files:
        return jsonify(*json_err('No file provided'))

    src, meta = save_upload(request.files['file'])
    if not src:
        return jsonify(*json_err(meta.get('error', 'Upload failed')))

    dest = unique_path('pptx')

    try:
        from pptx import Presentation
        from pptx.util import Emu
        from PIL import Image

        # Render pages at 150 dpi
        images_bytes, page_count = _pdf_to_images(src, 'png', dpi=150)

        # Determine slide dimensions from first image
        first_img = Image.open(io.BytesIO(images_bytes[0]))
        w_px, h_px = first_img.size
        # Convert px @ 150dpi → EMU (914400 EMU per inch)
        emu_w = int(w_px / 150 * 914400)
        emu_h = int(h_px / 150 * 914400)

        prs = Presentation()
        prs.slide_width  = Emu(emu_w)
        prs.slide_height = Emu(emu_h)

        blank_layout = prs.slide_layouts[6]  # completely blank

        for img_bytes in images_bytes:
            slide = prs.slides.add_slide(blank_layout)
            img_stream = io.BytesIO(img_bytes)
            slide.shapes.add_picture(img_stream, 0, 0, width=Emu(emu_w), height=Emu(emu_h))

        prs.save(dest)

    except ImportError as exc:
        return jsonify(*json_err(f'Missing library: {exc}. Run: pip install python-pptx Pillow'))
    except Exception as exc:
        return jsonify(*json_err(f'PowerPoint conversion failed: {exc}'))
    finally:
        try:
            os.remove(src)
        except OSError:
            pass

    return jsonify(**json_ok({
        'download_url': download_url(dest),
        'slide_count': page_count,
        'output_size': os.path.getsize(dest),
        'warning': meta.get('warning'),
    })[0])
